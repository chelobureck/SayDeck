import openai
from config.settings import get_settings
from typing import Dict, Any
from pptx import Presentation as PPTXPresentation
import tempfile

settings = get_settings()
openai.api_key = settings.OPENAI_API_KEY

async def transcribe_audio(audio_file_path: str) -> str:
    """Транскрибирует аудио файл с помощью Whisper API"""
    with open(audio_file_path, "rb") as audio_file:
        transcript = await openai.Audio.atranscribe("whisper-1", audio_file)
    return transcript["text"]

async def generate_presentation_structure(text: str) -> Dict[str, Any]:
    """Генерирует структуру презентации с помощью GPT-4"""
    prompt = f"""
    Создай структуру презентации на основе следующего текста. 
    Верни JSON с полями:
    - title: заголовок презентации
    - slides: массив слайдов, где каждый слайд имеет поля:
        - title: заголовок слайда
        - content: основной текст слайда
        - type: тип слайда (title, content, image, etc.)
    
    Текст: {text}
    """
    
    response = await openai.ChatCompletion.acreate(
        model="gpt-4",
        messages=[
            {"role": "system", "content": "Ты - эксперт по созданию презентаций."},
            {"role": "user", "content": prompt}
        ],
        temperature=0.7
    )
    
    return response.choices[0].message.content 

async def generate_presentation_pptx(content: dict) -> str:
    prs = PPTXPresentation()
    for slide in content.get("slides", []):
        slide_layout = prs.slide_layouts[1]  # обычный слайд
        pptx_slide = prs.slides.add_slide(slide_layout)
        pptx_slide.shapes.title.text = slide.get("title", "")
        pptx_slide.placeholders[1].text = slide.get("content", "")
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
    prs.save(tmp.name)
    return tmp.name 